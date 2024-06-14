
from django.shortcuts import render,redirect,get_object_or_404
from .models import UploadedFile
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth import logout as django_logout
import os
import io
import fitz  # PyMuPDF for PDF parsing
import docx  # python-docx for .docx parsing
import textract  # textract for general file parsing
import csv
from pptx import Presentation
from PIL import Image

from django.http import HttpResponseBadRequest
from django.http import HttpResponse
from .models import ChatSession,Feedback,SimilarQuestion

from sklearn.metrics.pairwise import cosine_similarity,euclidean_distances

from django.urls import reverse_lazy

from django.contrib.auth.views import LoginView

import numpy as np
import nltk
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
import inflect
nltk.download('punkt')
nltk.download('wordnet')
import re
from sentence_transformers import SentenceTransformer
from .forms import QnAForm
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from django.utils.crypto import get_random_string
from django.utils import timezone
import uuid 
import openai 
import fitz



nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')



from django.urls import reverse


def index(request):
    return render(request,'index.html')


 

@csrf_exempt
@login_required(login_url='/admin_login/')
def upload_file(request):
    if request.method == 'POST':
        if 'file' in request.FILES:
            uploaded_file = request.FILES['file']
            current_user = request.user
            chatbot_name = request.POST.get('bot_name')
            chatbot_description = request.POST.get('bot_desc')

            # Save the uploaded file to the database
            uploaded_file_obj = UploadedFile.objects.create(file=uploaded_file, user=current_user,chatbot_name=chatbot_name,chatbot_description=chatbot_description)

            # Get the ID of the uploaded file
            file_id = uploaded_file_obj.id

            # Redirect to the chat page after uploading the file, including the file ID in the URL
            return redirect(reverse('chat', kwargs={'user_name': current_user, 'file_id': file_id})+ '?success=true')
        elif 'link' in request.POST:
            link = request.POST['link']
            print(link)
            current_user = request.user
            chatbot_name = request.POST.get('bot_name')
            chatbot_description = request.POST.get('bot_desc')

            
            response = requests.get(link)
            if response.status_code == 200:
                # If the request is successful, parse the content
                soup = BeautifulSoup(response.content, 'html.parser')
                paragraphs = soup.find_all('p')
                content = '\n'.join([p.get_text() for p in paragraphs])
                
                # Save the link and content to the database
                uploaded_file_obj = UploadedFile.objects.create(link=link, content=content, user=current_user,chatbot_name=chatbot_name,chatbot_description=chatbot_description)
                file_id = uploaded_file_obj.id
            

            return redirect(reverse('chat', kwargs={'user_name': current_user, 'file_id': file_id})+ '?success=true')
            # else:
            #     # If the request fails, handle the error
            #     return HttpResponseBadRequest("Failed to fetch content from the provided URL.")

        else:
            # If neither a file nor a link is provided, return a bad request response
            return HttpResponseBadRequest("Invalid request. Please provide either a file or a link.")
    
    return render(request, 'upload_file.html')

from bs4 import BeautifulSoup
import requests




def clear_all_sessions(request):
    request.session.flush()
    return HttpResponse("All sessions have been cleared.")



# class CustomLoginView(LoginView):
#     template_name = "admin_login.html"  # Use your login template
#     redirect_authenticated_user = (
#         True  # Redirect logged-in users away from the login page
#     )

#     def get_success_url(self):
#         if self.request.user.is_authenticated and self.request.user.is_superuser:
#             return reverse_lazy("dashboard")
#         else:
#             return reverse_lazy("unauthorised")

#     def form_valid(self, form):
#         return super().form_valid(form)

#     def form_invalid(self, form):
#         return super().form_invalid(form)



from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth import authenticate, login
def custom_login_view(request):
    # Redirect authenticated users
    if request.user.is_authenticated:
        if request.user.is_superuser:
            return HttpResponseRedirect(reverse_lazy("dashboard"))
        else:
            return HttpResponseRedirect(reverse_lazy("unauthorised"))

    if request.method == "POST":
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            user = authenticate(username=form.cleaned_data['username'], password=form.cleaned_data['password'])
            if user is not None:
                login(request, user)
                return redirect(get_success_url(request))
            else:
                # Handle the case where authentication fails
                return render(request, "admin_login.html", {'form': form})
        else:
            # Form is invalid
            return render(request, "admin_login.html", {'form': form})
    else:
        # GET request, show the empty form
        form = AuthenticationForm()
        return render(request, "admin_login.html", {'form': form})



def get_success_url(request):
    if request.user.is_authenticated and request.user.is_superuser:
        return reverse_lazy("dashboard")
    else:
        return reverse_lazy("unauthorised")











@login_required(login_url='/admin_login/')
def dashboard_view(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    uploaded_files = UploadedFile.objects.filter(user=request.user)
    for uploaded_file in uploaded_files:
        uploaded_file.chat_url = request.build_absolute_uri(
            reverse('chat', kwargs={'user_name': request.user.username, 'file_id': uploaded_file.id})
        )
        # uploaded_file.chat_url = reverse('chat', kwargs={'user_name': request.user.username, 'file_id': uploaded_file.id})
        print(uploaded_file.chat_url)  # Debug print
    
    return render(request, 'dashboard.html', {'uploaded_files': uploaded_files})



@login_required(login_url='/admin_login/')
def chatbot_details(request, file_id):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    uploaded_file = get_object_or_404(UploadedFile, user=request.user, id=file_id)
    
    return render(request, 'details.html', {'uploaded_file': uploaded_file})




@login_required(login_url='/admin_login/')
def unique_questions_table(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
        
    
    user_entries = ChatSession.objects.filter(user=request.user).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    global_displayed_files = set()  # Global set to track displayed files
    
    for entry in user_entries:
        user_id = entry.user_identifier
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': [], 'displayed_files': set()}
        
        if entry.file_id not in global_displayed_files:  
            global_displayed_files.add(entry.file_id)  
            try:
                uploaded_file = UploadedFile.objects.get(pk=entry.file_id)
                entry.file_name = uploaded_file.file.name
                entry.chatbot_name = UploadedFile.chatbot_name
                answers_by_user[user_id]['answers'].append(entry)
            except UploadedFile.DoesNotExist:
                continue  

    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time,
    }
    
    return render(request, 'unique_questions.html', context)




@login_required(login_url='/admin_login/')
def file_queries(request, file_id):
    # uploaded_file = get_object_or_404(UploadedFile, pk=file_id)
    user_entries = ChatSession.objects.filter(user=request.user,file_id=file_id).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    
    for entry in user_entries:
        user_id = entry.user_identifier
        
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': []}
        
        answers_by_user[user_id]['answers'].append(entry)


        uploaded_file = get_object_or_404(UploadedFile, pk=entry.file_id)
        
        # Append the file name to the entry
        entry.file_name = uploaded_file.file.name

    
    
    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time
    }
    return render(request, 'file_details.html', context)

from django.utils import timezone


@login_required(login_url='/admin_login/')
def file_queries_log(request, file_id):
    # uploaded_file = get_object_or_404(UploadedFile, pk=file_id)
    user_entries = SimilarQuestion.objects.filter(user_name=request.user,file_id=file_id).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    
    for entry in user_entries:
        user_id = entry.user_session_id
        
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': []}
        
        answers_by_user[user_id]['answers'].append(entry)


        uploaded_file = get_object_or_404(UploadedFile, pk=entry.file_id)
        
        # Append the file name to the entry
        entry.file_name = uploaded_file.file.name

    
    
    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time
    }
    return render(request, 'file_details_log.html', context)









@login_required(login_url='/admin_login/')
def dashboard_log_question(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    user_entries = SimilarQuestion.objects.filter(user_name=request.user).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    global_displayed_files = set()  # Global set to track displayed files
    
    for entry in user_entries:
        user_id = entry.user_session_id
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': [], 'displayed_files': set()}
        
        if entry.file_id not in global_displayed_files:  # Check against the global set
            global_displayed_files.add(entry.file_id)  # Add to the global set
            try:
                uploaded_file = UploadedFile.objects.get(pk=entry.file_id)
                entry.file_name = uploaded_file.file.name
                answers_by_user[user_id]['answers'].append(entry)
            except UploadedFile.DoesNotExist:
                continue  # Skip adding entry if the file does not exist

    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time,
    }
    
    return render(request, 'log_questions.html', context)



def logout(request):
    django_logout(request)
    return redirect('admin_login')




def qna_create(request):
    if request.method == 'POST':
        form = QnAForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('dashboard')
    else:
        form = QnAForm()  # Initialize the form for GET requests
    return render(request, 'qna_form.html', {'form': form})



def qna_update(request, pk):
    qna = get_object_or_404(ChatSession, pk=pk)
    if request.method == 'POST':
        form = QnAForm(request.POST, instance=qna)
        if form.is_valid():
            form.save()
            return redirect('file_details')
    else:
        form = QnAForm(instance=qna)
    return render(request, 'qna_form.html', {'form': form, 'qna': qna}) 





def qna_delete(request, pk):
    qna = get_object_or_404(ChatSession, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('unique_questions')
    return render(request, 'qna_confirm_delete.html', {'qna': qna})

from django.contrib.auth.models import User

def qna_delete_log(request, pk):
    qna = get_object_or_404(SimilarQuestion, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('file_details_log_questions')
    return render(request, 'qna_confirm_delete.html', {'qna': qna})

def chatbot_delete(request, pk):
    qna = get_object_or_404(UploadedFile, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('dashboard')
    return render(request, 'chatbot.html', {'qna': qna})


@login_required(login_url='/admin_login/')
def verify_question(request, session_id, file_id):
    sessions = ChatSession.objects.filter(file_id=file_id, session_id=session_id)
    if not sessions:
        return JsonResponse({'message': 'No matching session found'}, status=404)
    
    verified = False
    for session in sessions:
        if request.user not in session.verified_by.all():
            session.verification_count += 1
            session.verified_by.add(request.user)
            session.save()
            verified = True
    
    if verified:
        return JsonResponse({'message': 'Verification successful'})
    else:
        return JsonResponse({'message': 'Already verified by this user'})
    


import os
import fitz  # PyMuPDF
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import json
import boto3





MAX_TOKENS_USAGE = 9000
# MAX_QUERIES_ALLOWED = 1
import PyPDF2
from django.conf import settings
import tempfile
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text()
    except Exception as e:
        print(f"Error occurred while extracting text from PDF: {e}")
        raise e
    return text

def split_text(text, max_length=8000):
    sentences = text.split('. ')
    segments = []
    current_segment = ""
    
    for sentence in sentences:
        if len(current_segment) + len(sentence) <= max_length:
            current_segment += sentence + '. '
        else:
            segments.append(current_segment.strip())
            current_segment = sentence + '. '
    
    if current_segment:
        segments.append(current_segment.strip())
    
    return segments

def find_relevant_segment(question, segments):
    vectorizer = TfidfVectorizer().fit_transform([question] + segments)
    vectors = vectorizer.toarray()
    
    question_vector = vectors[0]
    segment_vectors = vectors[1:]
    
    similarities = cosine_similarity([question_vector], segment_vectors)
    most_relevant_index = np.argmax(similarities)
    
    return segments[most_relevant_index]


import json
import random

# GREETINGS_INPUT = ["hi", "hello", "hey", "greetings", "good day", "good morning", "good afternoon", "good evening"]
# # Responses the system can give to greetings
# GREETINGS_RESPONSE = ["Hello!", "Hi there!", "Hello! How can I assist you today?","Good Morning", "Good Afternoon", "Good Evening"]

GREETINGS_MAP = {
    "hi": "Hi there!",
    "hello": "Hello! How can I assist you today?",
    "hey": "Hey! What can I do for you?",
    "greetings": "Greetings! How may I help?",
    "good day": "Good day to you!",
    "good morning": "Good morning! Hope you're having a great day!",
    "good afternoon": "Good afternoon! How can I assist?",
    "good evening": "Good evening! What do you need help with?",
    "how are you": "I'm just a bot, but thank you for asking! How can I assist you today?"
}

@csrf_exempt
def chat(request, user_name, file_id):
    print(user_name)
    file_exists = False

    if request.method == 'POST':
        # Handle POST request
        pass
    elif request.method == 'GET':
        file_exists = UploadedFile.objects.filter(id=file_id, user__username=user_name).exists()

    if request.method == 'GET' and file_exists:
        return render(request, 'chat.html', {
            'user_name': user_name,
            'file_id': file_id,
            'file_exists': file_exists,
        })
    elif request.method == 'GET' and not file_exists:
        return HttpResponse("Link Disabled.", status=404)
    elif request.method == 'POST':
        if 'user_identifier' not in request.session:
            user_identifier = f"ch-{uuid.uuid4().hex[:2]}"
            request.session['user_identifier'] = user_identifier
        else:
            user_identifier = request.session['user_identifier']

        session_start_time = request.session.get('session_start_time')
        if not session_start_time:
            session_start_time = timezone.localtime(timezone.now())
            request.session['session_start_time'] = session_start_time.strftime("%Y-%m-%d %H:%M:%S")

        session_id = f"{user_identifier}.{uuid.uuid4().hex[:2]}"

        question = request.POST.get('question', '')
        if not question:
            return HttpResponseBadRequest('Missing or empty "question" parameter in the request.')
        
        
        # if question in GREETINGS_INPUT:
        #     random_greeting = random.choice(GREETINGS_RESPONSE)  # Select a random greeting response
        #     question_asked_time = timezone.localtime(timezone.now())
        #     return JsonResponse({
        #         'answer': random_greeting,
        #         'source': 'system',
        #         'session_id': session_id,
        #         'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
                
                
        #     })



        if question in GREETINGS_MAP:
            response_message = GREETINGS_MAP[question]
            return JsonResponse({
                'answer': response_message,
                'source': 'system'
            })
        


        question_asked_time = timezone.localtime(timezone.now())
        request.session['current_time'] = question_asked_time.strftime("%Y-%m-%d %H:%M:%S")

        existing_answer = get_answer_from_database(question, user_identifier, session_id, user_name, file_id, question_asked_time)

        if existing_answer:
            answer, verification_count, positive_count, negative_feedback_count, file_id, question_asked_time = existing_answer
            if verification_count is not None:
                return JsonResponse({
                    'answer': answer,
                    'source': 'database',
                    'session_id': session_id,
                    'verification_count': verification_count,
                    'positive_count': positive_count,
                    'negative_feedback_count': negative_feedback_count,
                    'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
                })
            else:
                return JsonResponse({
                    'answer': answer,
                    'source': 'database',
                    'session_id': session_id
                })

        user = User.objects.get(username=user_name)
        uploaded_file = get_object_or_404(UploadedFile, id=file_id, user__username=user_name)

        if uploaded_file and uploaded_file.link:
            link = uploaded_file.link
            try:
                response = requests.get(link)
                if response.status_code == 200:
                    if link.endswith(".pdf"):
                        filename = link.split("/")[-1]
                        pdf_path = os.path.join(os.getcwd(), filename)

                        with open(pdf_path, 'wb') as pdf_file:
                            pdf_file.write(response.content)

                        document_text = extract_text_from_pdf(pdf_path)
                        os.remove(pdf_path)
                    else:
                        soup = BeautifulSoup(response.content, 'html.parser')
                        paragraphs = soup.find_all('p')
                        document_text = '\n'.join([p.get_text() for p in paragraphs])
                else:
                    return JsonResponse({'error': 'Failed to fetch content from the provided link.'})
            except Exception as e:
                return JsonResponse({'error': str(e)})
        elif uploaded_file and uploaded_file.file:
            file_path = uploaded_file.file.path
            file_extension = os.path.splitext(uploaded_file.file.name)[1].lower()

            if file_extension == '.pdf':
                document_text = extract_text_from_pdf(file_path)
            elif file_extension == '.docx':
                doc = docx.Document(file_path)
                document_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif file_extension == '.pptx':
                prs = Presentation(file_path)
                document_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif file_extension == '.txt':
                with open(file_path, 'r') as file:
                    document_text = file.read()
            elif file_extension == '.csv':
                with open(file_path, newline='', encoding='utf-8') as csvfile:
                    csvreader = csv.reader(csvfile)
                    document_text = "\n".join([",".join(row) for row in csvreader])
            elif file_extension == '.gif':
                image = Image.open(file_path)
                document_text = pytesseract.image_to_string(image)
            else:
                document_text = textract.process(file_path, method='text', encoding='utf-8')
        else:
            return JsonResponse({'error': 'No uploaded file found.'})

        text_segments = split_text(document_text)

        

        from transformers import GPT2Tokenizer

# Initialize the tokenizer
        tokenizer = GPT2Tokenizer.from_pretrained('gpt2')

        def count_tokens(text):
            tokens = tokenizer.encode(text, add_special_tokens=False)
            return len(tokens)

        def clean_response(response_text, stop_sequences=["<end>", "</s>", "</n>"]):
            # Find the earliest occurrence of any stop sequence
            first_stop_index = min((response_text.find(seq) if seq in response_text else len(response_text) for seq in stop_sequences))
            
            # Truncate the response at the first stop sequence
            response_text = response_text[:first_stop_index].strip()

            # Remove any known unwanted tokens
            unwanted_tokens = ["ASSISTANT:", "assistant:", "user:"]
            for token in unwanted_tokens:
                response_text = response_text.replace(token, "").strip()

            return response_text

        from transformers import GPT2Tokenizer

# Initialize the Bedrock client globally
        bedrock_client = boto3.client(service_name='bedrock-runtime')

        # Initialize the tokenizer
        tokenizer = GPT2Tokenizer.from_pretrained('gpt2')

       

        def initialize_bedrock_client():
            # Create a Bedrock client
            return boto3.client(service_name='bedrock-runtime')
        
        def truncate_context(context, max_tokens):
            tokens = tokenizer.encode(context, add_special_tokens=False)
            if len(tokens) > max_tokens:
                truncated_tokens = tokens[:max_tokens]
                return tokenizer.decode(truncated_tokens)
            return context

        def get_response(file_id, relevant_text,context, question, stop_sequences=["<end>", "</s>", "</n>"], max_attempts=3):
            print(relevant_text)
            max_context_tokens = 8192 - 1000  # Reserve tokens for the user input and response
            truncated_context = truncate_context(context, max_context_tokens)
            print(truncated_context)
            messages = [
                {"role": "system", "content": "You are a helpful assistant. Provide a brief, concise answer using only the uploaded document's content."},
                {"role": "system", "content": "You are a helpful assistant. Answer the following question based solely on the uploaded document content. Be brief and terminate your response with '<end>'."},
                {"role": "system", "content": "You are a helpful assistant. Always add a full stop at the end of the response"},
                {"role": "system", "content": "If the question is related to politicians ,film industry, cricketers,singers, languages and is outside the scope of document, respond with 'I don't know' and terminate your response with '<end>'."},
                {"role": "user", "content": f"Document details (File ID: {file_id}):"},
                {"role": "user", "content": truncated_context},
                
                {"role": "user", "content": question}
            ]
            full_prompt = "\n".join([f"{msg['role']}: {msg['content']}" for msg in messages]) + "\nassistant:"

            body = json.dumps({
                "prompt": full_prompt,
                "temperature": 0.1,
                "top_p": 0.9
            })

            modelId = 'meta.llama3-8b-instruct-v1:0'
            accept = 'application/json'
            contentType = 'application/json'

            # Invoke the model with the input body
            response = bedrock_client.invoke_model(body=body, modelId=modelId, accept=accept, contentType=contentType)

            # Read the StreamingBody object as a string
            response_body_str = response['body'].read().decode('utf-8')
            try:
                response_text = ""
                response_body = json.loads(response_body_str)
                # print("Debug: Response Body:", response_body)  # Debugging line
                
                # Extract and return the generation text if available
                generation = response_body.get('generation')
                print(f"Generation Text :{generation}")
                response_text += generation
                clean_resp = clean_response(response_text, stop_sequences)
                token_count = count_tokens(clean_resp)
                print(f"Token count of the response: {token_count}")
                # return generation if generation else "Sorry, I didn't understand that."
                return clean_response(generation) if generation else "Sorry, I didn't understand that."
            except json.JSONDecodeError as e:
                print(f"Error decoding JSON: {str(e)}")
                print("Debug: Response Body String:", response_body_str)  # Debugging line
                return "Sorry, there was an error processing the response."
            
        bedrock_client = initialize_bedrock_client()

            
        

            

        relevant_segment = find_relevant_segment(question, text_segments)
        # prompt = f"Context: {relevant_segment}\n\nQuestion: {question}\n### Answer:"
        # Assuming 'relevant_segment' is the context and 'question' is the actual question being asked.
        response = get_response(file_id,bedrock_client, relevant_segment, question)


        print(f"Chatbot: {response}")

        cluster_id = generate_cluster_id(question, user_identifier, session_start_time)
        print(cluster_id)

        save_to_database(question, response, user_identifier, session_id, session_start_time, cluster_id, user_name, file_id, question_asked_time)

        return JsonResponse({
            'answer': response,
            'source': 'llm',
            'session_id': session_id,
            'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
        })

    return HttpResponseBadRequest('Invalid request.')




import hashlib



def generate_cluster_id(question, user_identifier, session_start_time):
    # Concatenate question, user_identifier, and session_start_time
    concatenated_str = f"{question}{user_identifier}{session_start_time}"
    
    # Hash the concatenated string using SHA-256
    hash_object = hashlib.sha256(concatenated_str.encode())
    
    # Get the hexadecimal representation of the hash
    hex_digest = hash_object.hexdigest()
    
    # Take the first 5 characters of the hexadecimal representation
    shortened_hex_digest = hex_digest[:5]
    
    # Combine the prefix with the shortened hexadecimal representation
    cluster_id = f"CL-ID_{shortened_hex_digest}"
    
    return cluster_id





def preprocess_text(text):
    lemmatizer = WordNetLemmatizer()
    p = inflect.engine()

    def convert_number_to_words(token):
        if token.isdigit():
            return p.number_to_words(token)
        else:
            return token

    tokens = word_tokenize(text)
    lemmatized_tokens = [lemmatizer.lemmatize(convert_number_to_words(token.lower())) for token in tokens]
    # print('lemmatized_tokens',lemmatized_tokens)
    return ' '.join(lemmatized_tokens)







def extract_session_number(question):
    # Extracting session number from the question
    match = re.search(r'session\s*(\d+)', question, re.IGNORECASE)
    # print('match',match)
    return int(match.group(1)) if match else None





def get_answer_from_database(question,user_identifier,session_id,user_name,file_id,question_asked_time):
    try:
        all_entries = ChatSession.objects.filter(user__username=user_name,file_id=file_id)

        
        if not all_entries:
            return None  # No questions in the database, return None
        
        current_question_session_number = extract_session_number(question)
        processed_current_question = preprocess_text(question)
        # print('processed_current_questions',processed_current_question)

        # Filter entries by session number if present
        relevant_entries = [entry for entry in all_entries if extract_session_number(entry.question) == current_question_session_number]
        # print('relevant_entries',relevant_entries)

        if not relevant_entries:
            return None  # No relevant entries found for the specific session

        processed_questions = [preprocess_text(entry.question) for entry in relevant_entries]
        processed_questions.append(processed_current_question)  # Adding the current question for similarity comparison
        # print('processed_questions',processed_questions)

        # Use SentenceTransformer for creating embeddings
        model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
        question_embeddings = model.encode(processed_questions, convert_to_tensor=True)

        # Convert embeddings to NumPy arrays for similarity calculation
        question_embeddings_np = [embedding.cpu().detach().numpy() for embedding in question_embeddings]

        # Calculate cosine similarities
        similarities = cosine_similarity([question_embeddings_np[-1]], question_embeddings_np[:-1])
        # print('similarities',similarities)

        # Find the most similar question
        most_similar_index = similarities.argmax()
        most_similar_answer = relevant_entries[most_similar_index].answer
        verification_count = relevant_entries[most_similar_index].verification_count
        positive_count = relevant_entries[most_similar_index].positive_count
        negative_feedback_count = relevant_entries[most_similar_index].negative_feedback_count

        # Adjust similarity threshold based on context
        similarity_threshold = 0.8 if 'session' in processed_current_question else 0.8

        if similarities[0][most_similar_index] >= similarity_threshold:
            retrieved_entry = relevant_entries[most_similar_index]
            # if not SimilarQuestion.objects.filter(question=retrieved_entry.question, user_session_id=user_identifier).exists():
            retrieved_entry.retrieval_count += 1
            retrieved_entry.save() 
            # existing_similar_question = SimilarQuestion.objects.filter(question=retrieved_entry.question, user_session_id=user_identifier).exists()
            # if not existing_similar_question:
            
            # Create a SimilarQuestion object
            match = SimilarQuestion.objects.create(
                    question=question,
                    answer=retrieved_entry.answer,
                    session_id=session_id,
                    user_session_id=user_identifier,
                    cluster_id=retrieved_entry.cluster_id,
                    session_start_time=retrieved_entry.session_start_time,
                    question_asked_time=question_asked_time,
                    user_name=user_name,
                    file_id=file_id
                )

            return most_similar_answer,verification_count,positive_count,negative_feedback_count,file_id,question_asked_time
        else:
            return None

    except ChatSession.DoesNotExist:
        return None


def count_positive_feedback(question):
    try:
        # Count the number of positive feedback for the given question
        positive_feedback_count = Feedback.objects.filter(question=question, feedback=True).count()
        return positive_feedback_count
    except Feedback.DoesNotExist:
        return 0

def count_negative_feedback(question):
    try:
        # Count the number of positive feedback for the given question
        negative_feedback_count = Feedback.objects.filter(question=question, feedback=False).count()
        return negative_feedback_count
    except Feedback.DoesNotExist:
        return 0


import json
from django.db.models import F
from django.db.models import Max
from django.contrib.sessions.models import Session
from django.contrib.sessions.backends.db import SessionStore

@csrf_exempt
def submit_feedback(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body.decode('utf-8'))
            question = data.get('question')
            feedback = data.get('feedback')

            # Convert feedback to boolean
            feedback_bool = feedback.lower() == "positive"
            
            # Check if the user has already given feedback for this question

            # Create Feedback object
            feedback_obj = Feedback.objects.create(
                user_identifier=request.session.get('user_identifier'),
                session_id=data.get('session_id'),
                question=question,
                answer=data.get('answer'),
                feedback=feedback_bool
            )
            
            # Increase verification count if feedback is positive
            if feedback_bool:
                # Check if session data has processed feedbacks
                processed_feedbacks = request.session.get('processed_feedbacks', [])

                # Check if this feedback has already been processed
                if feedback_obj.id not in processed_feedbacks:
                    # Count positive feedback for this specific question
                    positive_feedback_count = Feedback.objects.filter(question=question, feedback=True, id__gt=feedback_obj.id).count()

                    # Update verification count for the question
                    latest_chat_session = ChatSession.objects.filter(question=question).latest('id')
                    latest_chat_session.positive_count += positive_feedback_count + 1  # Increment by 1 for the current feedback
                    latest_chat_session.save()

                    # Add this feedback to the processed list
                    processed_feedbacks.append(feedback_obj.id)
                    request.session['processed_feedbacks'] = processed_feedbacks

            else:  # If feedback is negative
                # Check if session data has processed negative feedbacks
                processed_negative_feedbacks = request.session.get('processed_negative_feedbacks', [])

                # Check if this negative feedback has already been processed
                if feedback_obj.id not in processed_negative_feedbacks:
                    # Count negative feedback for this specific question
                    negative_feedback_count = Feedback.objects.filter(question=question, feedback=False, id__gt=feedback_obj.id).count()

                    # Update negative feedback count for the question
                    latest_chat_session = ChatSession.objects.filter(question=question).latest('id')
                    latest_chat_session.negative_feedback_count += negative_feedback_count + 1  # Increment by 1 for the current negative feedback
                    latest_chat_session.save()

                    # Add this negative feedback to the processed list
                    processed_negative_feedbacks.append(feedback_obj.id)
                    request.session['processed_negative_feedbacks'] = processed_negative_feedbacks
                else:
                    # Negative feedback already processed, no need to count again
                    pass

            return JsonResponse({'status': 'success'})
        except json.JSONDecodeError:
            return JsonResponse({'status': 'error', 'message': 'Invalid JSON'}, status=400)
    else:
        return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)






def save_to_database(question, answer, user_identifier, session_id, session_start_time, cluster_id, user_name,file_id,question_asked_time):
   
    try:
        user = User.objects.get(username=user_name)
    except User.DoesNotExist:
        # Handle the case where the user does not exist
        return None  # or raise an exception

    # Create a ChatSession object
    chat_session = ChatSession.objects.create(
        session_id=session_id,
        session_start_time=session_start_time,
        question=question,
        question_asked_time=question_asked_time,
        answer=answer,
        user_identifier=user_identifier,
        cluster_id=cluster_id,
        user=user,
        file_id=file_id
    )

    # Create a SimilarQuestion object
    similar_question = SimilarQuestion.objects.create(
        session_id=session_id,
        session_start_time=session_start_time,
        question=question,
        question_asked_time=question_asked_time,
        answer=answer,
        user_session_id=user_identifier,
        cluster_id=cluster_id,
        user_name=user_name,
        file_id=file_id
    )

    
    chat_session.save()
    similar_question.save()








from datetime import datetime


from django.http import HttpResponseRedirect
import csv

def upload_csv(request):
    if request.method == "POST":
        csv_file = request.FILES['csv_file']
        decoded_file = csv_file.read().decode('utf-8').splitlines()
        reader = csv.DictReader(decoded_file)
        for row in reader:
            session_id = row['session_id']
            session_start_time = parse_datetime(row['session_start_time'])
            question = row['question']
            question_asked_time = parse_datetime(row['question_asked_time'])
            answer = row['answer']
            user_identifier = row['user_identifier']
            retrieval_count = int(row['retrieval_count'])
            verification_count = int(row['verification_count'])
            cluster_id = row.get('cluster_id', None)  # Handling if cluster_id is missing in some rows
            
            ChatSession.objects.create(
                session_id=session_id,
                session_start_time=session_start_time,
                question=question,
                question_asked_time=question_asked_time,
                answer=answer,
                user_identifier=user_identifier,
                retrieval_count=retrieval_count,
                verification_count=verification_count,
                cluster_id=cluster_id
            )
        return HttpResponseRedirect('/admin/app1/chatsession/')  # Redirect to admin page after upload
    return render(request, "upload_csv.html")  # Render the upload form if GET request

def parse_datetime(datetime_str):
    try:
        return datetime.fromisoformat(datetime_str)
    except ValueError:
        # Handle datetime string in alternative format if necessary
        # For example: '2024-04-11 04:24:37+00:00' to '2024-04-11 04:24:37'
        return datetime.strptime(datetime_str.split('+')[0], '%Y-%m-%d %H:%M:%S')





# import json
# @csrf_exempt
# def submit_feedback(request):
#     if request.method == "POST":
        
#         try:
#             data = json.loads(request.body.decode('utf-8'))  # Ensure decoding of the body
#             question = data.get('question')
#             answer = data.get('answer')
#             feedback = data.get('feedback')
#             session_id = data.get('session_id')
#             user_identifier = request.session.get('user_identifier')  # Use get() to avoid KeyError
#             print(session_id)
            
#             # Convert feedback to boolean
#             if feedback.lower() == "positive":
#                 feedback_bool = True
#             else:
#                 feedback_bool = False

#             print('question :', question)
#             print('answer :', answer)
#             print('feedback : ', feedback)
            
#             # Create Feedback object
#             Feedback.objects.create(
#                 user_identifier=user_identifier,
#                 session_id=session_id,
#                 question=question,
#                 answer=answer,
#                 feedback=feedback_bool
#             )
            
#             # Increase verification count if feedback is positive
#             if feedback_bool:
#                 try:
#                     chat_session = ChatSession.objects.get(session_id=session_id)
#                     chat_session.verification_count += 1
#                     chat_session.save()
#                 except ChatSession.DoesNotExist:
#                     pass  # Handle the case when the chat session doesn't exist
                
#             return JsonResponse({'status': 'success'})
#         except json.JSONDecodeError:
#             return JsonResponse({'status': 'error', 'message': 'Invalid JSON'}, status=400)
#     else:
#         return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)




































# import re

# def elaborate_question(question):
#     # Define patterns and corresponding replacements for question elaboration
#     elaborations = {
#         r'\b(please|kindly)\b': '',
#         r'\b(how can|what are|explain)\b': 'Please explain',
#         # Add more patterns and replacements as needed
#     }

#     # Apply each elaboration pattern and replacement
#     for pattern, replacement in elaborations.items():
#         question = re.sub(pattern, replacement, question, flags=re.IGNORECASE)

#     # Capitalize the first letter of the question and add a question mark if needed
#     question = question.strip().capitalize()
#     if not question.endswith('?'):
#         question += '?'

#     return question

# def construct_prompt(question):
#     # Construct a detailed prompt based on the question
#     prompt = f"Question: {question}\n\nContext: Please search the text from the file for information related to this question and provide the answer."
#     return prompt

# --------------------------------------------------------------------current used code for chatbot-------------------------------------------------------------




# @csrf_exempt
# def chat(request):
#     if request.method == 'POST':
#         question = request.POST.get('question', '')
#         print(f"Received question: {question}")

#         if not question:
#             return HttpResponseBadRequest('Missing or empty "question" parameter in the request.')
#         # elaborated_question = elaborate_question(question)
#         # print(elaborated_question)
#         uploaded_file_obj = None
#         extension = None
#         # Check if the question exists in the database
        
#         existing_answer = get_answer_from_database(question)

#         if existing_answer:
#             # If the question exists in the database, return the stored answer
#             response_data = {
#                 'answer': existing_answer,
#                 'source': 'database',
#             }
#             print(response_data)
#             return JsonResponse(response_data)
            
#         else:
#             # If the question is not in the database, proceed with language model logic
#             dir_path = os.path.dirname(os.path.realpath(__file__))
#             print(dir_path)


#             openai.api_key = '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'

#             # Path to your document file
#             document_path = os.path.join(dir_path, 'innohealth.pdf')
#             print(document_path)

#             # Read the content of your document file
#             with fitz.open(document_path) as doc:
#                 text = ""
#                 for page in doc:  # Iterate through each page
#                     text += page.get_text()
#                     print(text)

#             # Prompt the user to enter a question
#             user_question = question

#             try:
#                 # Making a call to the OpenAI API using the Chat Completion endpoint
#                 response = openai.ChatCompletion.create(
#                     model="gpt-3.5-turbo",  # Ensure using a compatible model for chat
#                     messages=[
#                         {"role": "system", "content": "You are a helpful assistant. provide the concise answers yet it complete.If you don't know the answer, just say that you don't know, don't try to make up an answer."},
#                         {"role": "user", "content": f"{text[:69000]}"},
                        
#                         {"role": "user", "content": user_question}
#                     ]
#                 )
#                 print("Answer:", response.choices[-1].message['content'].strip())  # Print the generated response
#             except Exception as e:
#                 print(f"An error occurred: {str(e)}")  # Handle any errors that occur during the API call

            
#         save_to_database(question, response.choices[0].message['content'].strip())

#             # Other response handling and session update code...

#         response_data = {
#                 'answer': response.choices[0].message['content'].strip(),
#                 'source': 'llm',
#             }
        
#         print(response_data)
#         return JsonResponse(response_data)
#     return render(request, 'chat.html', {'error': 'Invalid request method.'})




























# ------------------------------------------------------------------//finish code for chatbot//------------------------------------------------------------------------





# import openai


# api_key = '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'

# def split_text_by_tokens(text, max_tokens):
#     # Tokenize the text
#     tokens = text.split()

#     # Split the tokens into chunks of max_tokens
#     token_chunks = [tokens[i:i+max_tokens] for i in range(0, len(tokens), max_tokens)]

#     # Join the token chunks into strings
#     text_chunks = [' '.join(chunk) for chunk in token_chunks]

#     return text_chunks

# @csrf_exempt
# def chat(request):
#     if request.method == 'POST':
#         question = request.POST.get('question', '')
#         print(f"Received question: {question}")

#         if not question:
#             return HttpResponseBadRequest('Missing or empty "question" parameter in the request.')
        
#         # Load the contents of the uploaded file
#         uploaded_file_obj = UploadedFile.objects.last()

#         if not uploaded_file_obj:
#             return render(request, 'chat.html', {'error': 'No uploaded file found.'})

#         filename, extension = os.path.splitext(uploaded_file_obj.file.name)

#         if extension == ".pdf":
#             loaders = [PyPDFLoader(uploaded_file_obj.file.path)]
#         elif extension == ".docx":
#             loaders = [Docx2txtLoader(uploaded_file_obj.file.path)]
#         else:
#             return render(request, 'chat.html', {'error': 'Invalid extension. Please upload pdf or docx files.'})

#         file_contents = []
#         for loader in loaders:
#             docs = loader.load()
#             for doc in docs:
#                 file_contents.append(doc.page_content)  # Example method name (replace with the actual method)
#                 # print(file_contents)

#         # Construct a prompt using the contents of the file and the user's question
#         file_text = "\n".join(file_contents)
#         # print("file_text",file_text)
#         # Split file_text into chunks
#         max_tokens = 8192  # Maximum tokens allowed by OpenAI model
#         text_chunks = split_text_by_tokens(file_text, max_tokens)

# # Initialize response_text before the for loop
#         response_text = ''

#         # Generate completions for each text chunk and collect responses
#         for chunk in text_chunks:
#             # Keep track of remaining tokens
#             remaining_tokens = len(chunk.split())
#             # Initialize an empty list to store completions for this chunk
#             chunk_responses = []
            
#             # Iterate until all tokens in the chunk have been processed
#             while remaining_tokens > 0:
#                 # Determine the number of tokens to send to the API in this iteration
#                 tokens_to_send = min(remaining_tokens, max_tokens)
#                 # Prepare the message content for this iteration
#                 chunk_message = " ".join(chunk_responses + [chunk.split()[0]])  # Include previous responses and start of current chunk
#                 print('Chunk Message : ',chunk_message)
#                 # Send request to OpenAI
#                 completion = openai.ChatCompletion.create(
#                 model="gpt-4",
#                 messages=[
#                     {
#                         "role": "system",
#                         "content": "Please generate a response using the detailed information provided in the next message."
#                     },
#                     {
#                         "role": "system",
#                         "content": chunk_message 
#                         # print(content) # Detailed context as the first user message
#                     },
#                     {
#                         "role": "user",
#                         "content": question  # Actual question as the second user message
#                     }
#                 ],
#                 max_tokens=tokens_to_send
#             )
#                 # max_tokens=tokens_to_send


#                 # Append the completion to chunk_responses
#                 chunk_responses.append(completion.choices[0].message.content.strip())
#                 # Update remaining tokens
#                 remaining_tokens -= tokens_to_send
            
#             # Append the responses for this chunk to the overall response_text
#             response_text += '\n'.join(chunk_responses) + '\n'

#         print(response_text)

#         # Save the question-answer pair in the database
#         save_to_database(question, response_text)

#         response_data = {
#             'answer': response_text,
#             'source': 'llm',
#         }
        
#         print(response_data)
#         return JsonResponse(response_data)
#     return render(request, 'chat.html', {'error': 'Invalid request method.'})
























#---------------------------------------------------------//code without retrieval--------------------------------------------
# def reset_session_timeout(request):
#     # This view is responsible for resetting the session timeout on user activity
#     request.session.modified = True
#     return JsonResponse({'status': 'success'})




# tokenizer = AutoTokenizer.from_pretrained("distilbert-base-cased-distilled-squad")
# model = AutoModelForQuestionAnswering.from_pretrained("distilbert-base-cased-distilled-squad")


# chat_history_key = 'chat_history'

# @csrf_exempt
# def chat(request):
#     if request.method == 'POST':
#         question = request.POST.get('question', '')
#         print(f"Received question: {question}")

        # if not question:
        #     return HttpResponseBadRequest('Missing or empty "question" parameter in the request.')

        # uploaded_file_obj = UploadedFile.objects.last()

        # if not uploaded_file_obj:
        #     return render(request, 'chat.html', {'error': 'No uploaded file found.'})

        # filename, extension = os.path.splitext(uploaded_file_obj.file.name)

        # if extension == ".pdf":
        #     loaders = [PyPDFLoader(uploaded_file_obj.file.path)]
        # elif extension == ".docx":
        #     loaders = [Docx2txtLoader(uploaded_file_obj.file.path)]
        # else:
        #     return render(request, 'chat.html', {'error': 'Invalid extension. Please upload pdf or docx files.'})

        # docs = []
        # for loader in loaders:
        #     docs.extend(loader.load())
        #     # print("Number of documents loaded:", len(docs))

        #     # # Assuming 'page_content' is the correct attribute to access the content
        #     # if hasattr(docs[0], 'page_content'):
        #     #     print("First document content:", getattr(docs[0], 'page_content'))
                
        #     # else:
        #     #     print("Content attribute not found in the Document object")

        # text_splitter = RecursiveCharacterTextSplitter(
        #     chunk_size=1500,
        #     chunk_overlap=50
        # )
        # splits = text_splitter.split_documents(docs)
        # # print("Number of splits:", len(splits))
        # # for i, split in enumerate(splits):
            
        # #     if hasattr(split, 'page_content'):
        # #         print(f"Content of split {i + 1}:", getattr(split, 'page_content'))
        # #     else:
        # #         print(f"Content attribute not found in Split {i + 1} Document object")

        # # # Assuming 'page_content' is the correct attribute to access the content
        # # if hasattr(splits[0], 'page_content'):
        # #     print("First split content:", getattr(splits[0], 'page_content'))
           
        # # else:
        # #     print("Content attribute not found in the Document object") 
        # persist_directory = 'chroma/'
        # vectordb = Chroma.from_documents(
        #     documents=splits,
        #     embedding=OpenAIEmbeddings(),
        #     persist_directory=persist_directory
        # )

#         llm = ChatOpenAI(model_name=llm_name, temperature=0)
#         qa_chain = ConversationalRetrievalChain.from_llm(llm, vectordb.as_retriever())

#         user_identifier = request.session.session_key
#         if user_identifier is None:
#             request.session.cycle_key()
#             user_identifier = request.session.session_key

#         chat_history = request.session.get(chat_history_key, [])

#         if 'is_new_user' in request.session:
#             if request.session['is_new_user']:
#                 chat_history.clear()
#             request.session['is_new_user'] = False

#         chat_history = [tuple(item) for item in chat_history]

#         response = qa_chain({"question": question, "chat_history": chat_history})
         
#         chat_history.append((question, response['answer']))

#         request.session[chat_history_key] = chat_history

#         if 'is_new_user' not in request.session:
#             request.session.cycle_key()
#             request.session['is_new_user'] = True

#         ChatSession.objects.create(user_id=user_identifier, question=question, answer=response['answer'])

#         user_chat_history = ChatSession.objects.filter(user_id=user_identifier)

#         # return render(request, 'chat.html', {'question': question, 'response': response['answer'], 'user_chat_history': user_chat_history})
#         response_data = {
#             'answer': response['answer'],
#         }
#         print(response_data)
#         return JsonResponse(response_data)
#     return render(request, 'chat.html', {'error': 'Invalid request method.'})

#---------------------------------------------------------//code without retrieval//-------------------------------------------------------------------------



# def preprocess_text(text):
#     lemmatizer = WordNetLemmatizer()
#     tokens = word_tokenize(text)
#     lemmatized_tokens = [lemmatizer.lemmatize(token.lower()) for token in tokens]
#     return ' '.join(lemmatized_tokens)



# def preprocess_text(text):
#     lemmatizer = WordNetLemmatizer()
#     p = inflect.engine()

#     def convert_number_to_words(token):
#         if token.isdigit():
#             return p.number_to_words(token)
#         else:
#             return token

#     tokens = word_tokenize(text)
#     lemmatized_tokens = [lemmatizer.lemmatize(convert_number_to_words(token.lower())) for token in tokens]
    
#     return ' '.join(lemmatized_tokens)




# def elaborate_question(question):
#     """
#     Function to elaborate the question into a more meaningful sentence.
#     This helps in distinguishing between different sessions.
#     """
#     if 'session' in question:
#         session_number = [int(s) for s in question.split() if s.isdigit()]
#         if session_number:
#             return f"Please provide details about the speakers in session number {session_number[0]}."
#     return question

# def get_answer_from_database(question):
#     try:
#         # Try to find a direct match for session numbers first
#         session_number = [int(s) for s in question.split() if s.isdigit()]
#         if session_number:
#             session_query = f"session {session_number[0]}"
#             for entry in ChatSession.objects.all():
#                 if session_query in entry.question:
#                     return entry.answer

#         # If no direct match, proceed with NLP-based similarity search
#         all_entries = list(ChatSession.objects.all())
        
#         if not all_entries:
#             return None

#         processed_question = preprocess_text(elaborate_question(question))
#         processed_questions = [preprocess_text(entry.question) for entry in all_entries]

#         processed_questions.append(processed_question)

#         model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
#         question_embeddings = model.encode(processed_questions, convert_to_tensor=True)

#         question_embeddings_np = np.array([embedding.cpu().numpy() for embedding in question_embeddings])

#         similarities = cosine_similarity(question_embeddings_np[-1].reshape(1, -1), question_embeddings_np[:-1])

#         most_similar_index = similarities.argmax()

#         similarity_threshold = 0.8 if 'session' in processed_question else 0.4

#         if similarities[0][most_similar_index] >= similarity_threshold:
#             return all_entries[most_similar_index].answer
#         else:
#             return None

#     except ChatSession.DoesNotExist:
#         return None











# 3rd correct code for sentence transformers
# def elaborate_question(question):
#     """
#     Function to elaborate the question into a more meaningful sentence.
#     This helps in distinguishing between different sessions.
#     """
#     if 'session' in question:
#         session_number = [int(s) for s in question.split() if s.isdigit()]
#         print('session number',session_number)
#         if session_number:
#             return f"Please provide details about the speakers in session number {session_number[0]}."
#         print('question',question)
#     return question

# def get_answer_from_database(question):
#     try:
#         all_entries = ChatSession.objects.all()
        
#         if not all_entries:
#             return None

#         processed_question = preprocess_text(elaborate_question(question))
#         print('processed_question',processed_question)
#         processed_questions = [preprocess_text(entry.question) for entry in all_entries]

#         processed_questions.append(processed_question)

#         model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
#         question_embeddings = model.encode(processed_questions, convert_to_tensor=True)

#         question_embeddings_np = [embedding.cpu().detach().numpy() for embedding in question_embeddings]
#         # print('question_embeddings_np',question_embeddings_np)

#         similarities = cosine_similarity(np.array([question_embeddings_np[-1]]), np.array(question_embeddings_np[:-1]))

#         most_similar_index = similarities.argmax()

#         if 'session' in processed_question:
#             similarity_threshold = 0.8
#         else:
#             similarity_threshold = 0.4

#         if similarities[0][most_similar_index] >= similarity_threshold:
#             return all_entries[int(most_similar_index)].answer
#         else:
#             return None

#     except ChatSession.DoesNotExist:
#         return None



























# 2nd correct code for sentence transformers
# def get_answer_from_database(question):
#     try:
#         all_entries = ChatSession.objects.all()
        
#         if not all_entries:
#             return None  # No questions in the database, return None

#         # Assuming ChatSession has 'question' and 'answer' fields
#         processed_questions = [preprocess_text(entry.question) for entry in all_entries]
#         processed_current_question = preprocess_text(question)
        
#         # Adding the current question to the list for similarity comparison
#         processed_questions.append(processed_current_question)
#         print("Processed Questions:", processed_questions)

#         # Use SentenceTransformer for creating embeddings
#         model = SentenceTransformer('paraphrase-MiniLM-L6-v2')  # You can choose a different model
#         question_embeddings = model.encode(processed_questions, convert_to_tensor=True)

#         # Convert PyTorch tensors to NumPy arrays
#         question_embeddings_np = [embedding.cpu().detach().numpy() for embedding in question_embeddings]
#         print('question_embeddings', question_embeddings)

#         # Calculate cosine similarities between the current question and all stored questions
#         similarities = cosine_similarity(np.array([question_embeddings_np[-1]]), np.array(question_embeddings_np[:-1]))
#         print('similarities', similarities)

#         # Find the index of the most similar question
#         most_similar_index = similarities.argmax()

#         # Retrieve the answer corresponding to the most similar question
#         most_similar_answer = all_entries[int(most_similar_index)].answer

#         # Set a threshold for similarity
#         if 'session' in processed_current_question or 'speaker' in processed_current_question:
#             similarity_threshold = 0.8
#         else:
#             similarity_threshold = 0.4  # Default threshold for other questions

#         print(f"Similarity Threshold for the current question: {similarity_threshold}")
#         print(f"Similarity Score for the current question: {similarities[0][most_similar_index]}")

#         if similarities[0][most_similar_index] >= similarity_threshold:
#             return most_similar_answer
#         else:
#             return None

#     except ChatSession.DoesNotExist:
#         return None







#correct one
# def get_answer_from_database(question):
#     try:
#         all_entries = ChatSession.objects.all()
        
#         if not all_entries:
#             return None  # No questions in the database, return None

#         # Assuming ChatSession has 'question' and 'answer' fields
#         processed_questions = [preprocess_text(entry.question) for entry in all_entries]
#         processed_current_question = preprocess_text(question)
        
#         # Adding the current question to the list for similarity comparison
#         processed_questions.append(processed_current_question)
#         print("Processed Questions:", processed_questions)

#         # Vectorize the questions using TF-IDF
#         vectorizer = TfidfVectorizer()
#         question_vectors = vectorizer.fit_transform(processed_questions)

#         # Calculate cosine similarities between the current question and all stored questions
#         similarities = cosine_similarity(question_vectors[-1], question_vectors[:-1])
#         print('similarities', similarities)

#         # Find the index of the most similar question
#         most_similar_index = similarities.argmax()

#         # Retrieve the answer corresponding to the most similar question
#         most_similar_answer = all_entries[int(most_similar_index)].answer

#         # Set a threshold for similarity
#         if 'session' in processed_current_question or 'speaker' in processed_current_question:
#             similarity_threshold = 0.7
#         else:
#             similarity_threshold = 0.4  # Default threshold for other questions

#         print(f"Similarity Threshold for current question: {similarity_threshold}")
#         print(f"Similarity Score for current question: {similarities[0][most_similar_index]}")

#         if similarities[0][most_similar_index] >= similarity_threshold:
#             return most_similar_answer
#         else:
#             return None

#     except ChatSession.DoesNotExist:
#         return None






































